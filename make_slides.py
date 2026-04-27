"""Generate VibeFinder AI final presentation as a .pptx file."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Brand colors ────────────────────────────────────────────────────────────
DARK_BG    = RGBColor(0x0F, 0x0F, 0x1A)   # near-black navy
ACCENT     = RGBColor(0x7C, 0x3A, 0xED)   # violet
ACCENT2    = RGBColor(0x06, 0xB6, 0xD4)   # cyan
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xDD)
GOLD       = RGBColor(0xF5, 0x9E, 0x0B)
GREEN      = RGBColor(0x10, 0xB9, 0x81)
RED        = RGBColor(0xEF, 0x44, 0x44)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


def new_prs() -> Presentation:
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank_slide(prs: Presentation):
    layout = prs.slide_layouts[6]   # completely blank
    return prs.slides.add_slide(layout)


def bg(slide, color: RGBColor = DARK_BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, fill_color, alpha=None):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    return shape


def add_text(slide, text, left, top, width, height,
             font_size=24, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox


def add_para(tf, text, font_size=18, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, italic=False, space_before=0):
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(space_before)
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return p


def accent_bar(slide, y=Inches(0.08), color=ACCENT):
    add_rect(slide, 0, y, SLIDE_W, Inches(0.06), color)


# ── Slide builders ───────────────────────────────────────────────────────────

def slide_title(prs):
    """Slide 1 — Title"""
    s = blank_slide(prs)
    bg(s)
    # gradient-ish side panel
    add_rect(s, 0, 0, Inches(4.8), SLIDE_H, RGBColor(0x1A, 0x0A, 0x3A))
    accent_bar(s, y=0, color=ACCENT)
    accent_bar(s, y=SLIDE_H - Inches(0.06), color=ACCENT2)

    # Big emoji / icon area
    add_text(s, "🎵", Inches(1.2), Inches(1.0), Inches(2.4), Inches(1.6),
             font_size=72, align=PP_ALIGN.CENTER)

    # Main title
    add_text(s, "VibeFinder AI",
             Inches(5.2), Inches(0.7), Inches(7.8), Inches(1.4),
             font_size=52, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

    add_text(s, "Music Recommendations Powered by Claude",
             Inches(5.2), Inches(2.0), Inches(7.8), Inches(0.7),
             font_size=24, color=ACCENT2, align=PP_ALIGN.LEFT)

    # Divider line
    add_rect(s, Inches(5.2), Inches(2.85), Inches(7.5), Inches(0.04), ACCENT)

    # Subtitle block
    txb = slide.shapes.add_textbox if False else \
          s.shapes.add_textbox(Inches(5.2), Inches(3.05), Inches(7.8), Inches(2.8))
    tf = txb.text_frame
    tf.word_wrap = True
    for line, sz, bold, col in [
        ("Final Project — AI 110",                       20, False, LIGHT_GRAY),
        ("",                                             10, False, WHITE),
        ("GitHub: github.com/guka199/applied-music-recommendations", 17, False, ACCENT2),
        ("",                                             10, False, WHITE),
        ("Guram Janashia  ·  April 2026",                18, False, LIGHT_GRAY),
    ]:
        add_para(tf, line, font_size=sz, bold=bold, color=col)

    # Left panel tag lines
    for y, txt, sz, col in [
        (Inches(3.0), "NL → Claude → Music", 16, ACCENT2),
        (Inches(3.5), "7 / 7 Tests Pass",    16, GREEN),
        (Inches(4.0), "Confidence Scored",   16, GOLD),
        (Inches(4.5), "Fallback Safe",       16, LIGHT_GRAY),
    ]:
        add_text(s, txt, Inches(0.3), y, Inches(4.2), Inches(0.5),
                 font_size=sz, color=col, align=PP_ALIGN.CENTER)


def slide_origin(prs):
    """Slide 2 — Where It Started"""
    s = blank_slide(prs)
    bg(s)
    accent_bar(s, color=ACCENT)

    add_text(s, "Where It Started — Modules 1–3",
             Inches(0.5), Inches(0.2), Inches(12), Inches(0.7),
             font_size=34, bold=True, color=WHITE)

    add_text(s, "Music Recommender Simulation",
             Inches(0.5), Inches(0.95), Inches(12), Inches(0.5),
             font_size=22, color=ACCENT2, italic=True)

    # Two column cards
    for col_x, title, items, accent_col in [
        (Inches(0.4), "What it did", [
            "Scored 18 songs against hand-crafted user profiles",
            "5 weighted features: genre, mood, energy, valence, tempo",
            "Sorted by score → returned top-5 recommendations",
            "Explained every score in plain English",
        ], ACCENT),
        (Inches(6.9), "What it proved", [
            "Simple math + right features = plausible results",
            "Genre bonus (+2.0) can override better emotional matches",
            "Catalog skew disadvantages niche genres (1 blues song)",
            '"Right by formula" ≠ "right for the human" — alignment gap',
        ], ACCENT2),
    ]:
        add_rect(s, col_x, Inches(1.6), Inches(6.0), Inches(0.45), accent_col)
        add_text(s, title,
                 col_x + Inches(0.15), Inches(1.62), Inches(5.7), Inches(0.4),
                 font_size=18, bold=True, color=DARK_BG)

        card = s.shapes.add_textbox(col_x, Inches(2.1), Inches(6.0), Inches(3.8))
        tf = card.text_frame
        tf.word_wrap = True
        for item in items:
            p = tf.add_paragraph()
            p.space_before = Pt(6)
            run = p.add_run()
            run.text = "  •  " + item
            run.font.size = Pt(17)
            run.font.color.rgb = LIGHT_GRAY

    # Bottom callout
    add_rect(s, Inches(0.4), Inches(6.1), Inches(12.5), Inches(1.0),
             RGBColor(0x1A, 0x1A, 0x30))
    add_text(s,
             "Key insight: recommenders are weighted scorecards. "
             "The sophistication is in choosing the right features — not the math.",
             Inches(0.6), Inches(6.18), Inches(12.1), Inches(0.8),
             font_size=17, italic=True, color=GOLD, align=PP_ALIGN.CENTER)


def slide_what_i_built(prs):
    """Slide 3 — What I Built"""
    s = blank_slide(prs)
    bg(s)
    accent_bar(s, color=ACCENT2)

    add_text(s, "What I Built — VibeFinder AI",
             Inches(0.5), Inches(0.2), Inches(12), Inches(0.7),
             font_size=34, bold=True, color=WHITE)

    add_text(s, "Extended the scorecard with a Claude-powered two-stage pipeline",
             Inches(0.5), Inches(0.95), Inches(12), Inches(0.5),
             font_size=20, color=LIGHT_GRAY, italic=True)

    stages = [
        ("Stage 1", "NL Parser", ACCENT,
         'python -m src.main --mode ai\n  --query "something chill\n  for late-night studying"',
         "↓  claude-opus-4-7\n↓  JSON schema output\n↓  {genre, mood, energy,\n   valence, tempo_bpm}"),
        ("Stage 2", "Weighted Scorer", RGBColor(0x06, 0x7A, 0x99),
         "Runs on parsed profile\nNot on raw text",
         "Scores all 18 songs\nReturns top-10 candidates\nExplains every score"),
        ("Stage 3", "AI Re-ranker", ACCENT,
         "claude-opus-4-7\nMusical reasoning\nConfidence scores 0–1",
         "Guardrail: filters\nhallucinated titles\nReturns top-5 + confidence"),
    ]

    arrow_x = [Inches(0.3), Inches(4.55), Inches(8.8)]
    arrow_w = Inches(4.0)

    for i, (label, title, col, left_txt, right_txt) in enumerate(stages):
        x = arrow_x[i]
        # Card background
        add_rect(s, x, Inches(1.7), arrow_w, Inches(4.5), RGBColor(0x18, 0x18, 0x2E))
        # Header bar
        add_rect(s, x, Inches(1.7), arrow_w, Inches(0.55), col)
        add_text(s, f"{label}  ·  {title}",
                 x + Inches(0.1), Inches(1.72), arrow_w - Inches(0.2), Inches(0.5),
                 font_size=17, bold=True, color=DARK_BG)
        # Content
        add_text(s, left_txt,
                 x + Inches(0.15), Inches(2.35), arrow_w - Inches(0.3), Inches(1.9),
                 font_size=14, color=ACCENT2)
        add_text(s, right_txt,
                 x + Inches(0.15), Inches(4.3), arrow_w - Inches(0.3), Inches(1.7),
                 font_size=14, color=LIGHT_GRAY)

        # Arrow between cards
        if i < 2:
            add_text(s, "→",
                     x + arrow_w + Inches(0.05), Inches(3.4), Inches(0.4), Inches(0.5),
                     font_size=28, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

    # Fallback note
    add_rect(s, Inches(0.3), Inches(6.3), Inches(12.7), Inches(0.85),
             RGBColor(0x1A, 0x1A, 0x2E))
    add_text(s,
             "Fallback path: any anthropic.APIError or JSON error → keyword heuristic profile → weighted scorer directly  [labeled in output]",
             Inches(0.5), Inches(6.38), Inches(12.3), Inches(0.7),
             font_size=15, color=GOLD, align=PP_ALIGN.CENTER)


def slide_architecture(prs):
    """Slide 4 — System Diagram"""
    s = blank_slide(prs)
    bg(s)
    accent_bar(s, color=ACCENT)

    add_text(s, "System Architecture",
             Inches(0.5), Inches(0.2), Inches(12), Inches(0.7),
             font_size=34, bold=True, color=WHITE)

    code = """\
USER: "something chill for late-night studying"
         │
         ▼  AI MODE
┌─────────────────────────────────────┐
│  Stage 1: NL Parser                 │  src/ai_recommender.py
│  claude-opus-4-7 + JSON schema      │  cache_control: ephemeral
│  → {genre:lofi, mood:chill,         │
│     energy:0.35, valence:0.60,      │
│     tempo_bpm:75}                   │
└──────────────┬──────────────────────┘
               │ structured profile
               ▼
┌─────────────────────────────────────┐
│  Stage 2: Weighted Scorer           │  src/recommender.py
│  score = genre×2.0 + mood×1.0       │  (unchanged from Modules 1-3)
│        + energy×2.0 + valence×1.5   │
│        + tempo×1.0                  │
│  → top-10 candidates                │
└──────────────┬──────────────────────┘
               │ candidates list
               ▼
┌─────────────────────────────────────┐
│  Stage 3: AI Re-ranker              │  Guardrail: only catalog
│  claude-opus-4-7 + JSON schema      │  titles pass through
│  Musical domain reasoning           │  + confidence scores 0–1
│  → top-5 re-ranked + confidence     │
└──────────────┬──────────────────────┘
               │
               ▼
         CLI OUTPUT"""

    add_rect(s, Inches(0.4), Inches(1.0), Inches(7.6), Inches(6.1),
             RGBColor(0x0A, 0x0A, 0x18))
    add_text(s, code,
             Inches(0.55), Inches(1.08), Inches(7.3), Inches(6.0),
             font_size=11, color=ACCENT2)

    # Right side: key design principles
    add_text(s, "Key Design Principles",
             Inches(8.4), Inches(1.0), Inches(4.6), Inches(0.5),
             font_size=20, bold=True, color=WHITE)

    principles = [
        (ACCENT,  "LLM as Ranker",
                  "Claude only re-orders songs that exist in the catalog — never invents them."),
        (GREEN,   "Guardrail",
                  "Hallucinated titles are filtered before output; omitted songs appended as safety net."),
        (GOLD,    "Fallback",
                  "Any API failure → weighted scorer runs directly. 100% uptime guaranteed."),
        (ACCENT2, "Explainability",
                  "Every score is broken down by feature. Users see exactly why a song ranked."),
        (LIGHT_GRAY, "Logging",
                  "Every Claude call written to logs/YYYY-MM-DD.log at DEBUG level for audit."),
    ]

    y = Inches(1.6)
    for col, heading, body in principles:
        add_rect(s, Inches(8.2), y, Inches(0.08), Inches(0.75), col)
        add_text(s, heading,
                 Inches(8.45), y, Inches(4.5), Inches(0.35),
                 font_size=15, bold=True, color=col)
        add_text(s, body,
                 Inches(8.45), y + Inches(0.32), Inches(4.5), Inches(0.5),
                 font_size=13, color=LIGHT_GRAY)
        y += Inches(0.95)


def slide_demo(prs):
    """Slide 5 — Live Demo / Sample Interactions"""
    s = blank_slide(prs)
    bg(s)
    accent_bar(s, color=GREEN)

    add_text(s, "Demo — Sample Interactions",
             Inches(0.5), Inches(0.2), Inches(12), Inches(0.7),
             font_size=34, bold=True, color=WHITE)

    demos = [
        ("Classic Mode — Chill Lofi Studier",
         'genre=lofi  mood=chill  energy=0.38',
         ACCENT,
         [("1. Library Rain  —  Paper Lanterns",  "Score: 7.41", "genre + mood + energy match"),
          ("2. Midnight Coding  —  LoRoom",        "Score: 7.36", "genre + mood + energy match"),
          ("3. Focus Flow  —  LoRoom",             "Score: 6.44", "genre + energy match")]),
        ('AI Mode — "melancholy rainy day, slow"',
         'Claude parsed → genre=blues  mood=melancholy  energy=0.28',
         ACCENT2,
         [("1. Crossroads Lament  —  Blue Delta",     "Score: 5.84  [AI confidence: 0.94]", "genre + mood + energy"),
          ("2. Library Rain  —  Paper Lanterns",       "Score: 4.21  [AI confidence: 0.72]", "energy + valence"),
          ("3. Spacewalk Thoughts  —  Orbit Bloom",    "Score: 3.99  [AI confidence: 0.65]", "energy + tempo")]),
    ]

    col_x = [Inches(0.3), Inches(6.7)]
    for i, (title, sub, col, results) in enumerate(demos):
        x = col_x[i]
        w = Inches(6.0)
        add_rect(s, x, Inches(1.1), w, Inches(0.48), col)
        add_text(s, title, x + Inches(0.12), Inches(1.13), w, Inches(0.4),
                 font_size=15, bold=True, color=DARK_BG)

        add_text(s, sub, x, Inches(1.65), w, Inches(0.4),
                 font_size=13, italic=True, color=LIGHT_GRAY)

        y = Inches(2.15)
        for rank_title, score, why in results:
            add_rect(s, x, y, w, Inches(1.25), RGBColor(0x14, 0x14, 0x28))
            add_text(s, rank_title, x + Inches(0.15), y + Inches(0.05),
                     w - Inches(0.3), Inches(0.4), font_size=14, bold=True, color=WHITE)
            add_text(s, score, x + Inches(0.15), y + Inches(0.42),
                     w - Inches(0.3), Inches(0.35), font_size=13, color=col)
            add_text(s, "Why: " + why, x + Inches(0.15), y + Inches(0.76),
                     w - Inches(0.3), Inches(0.35), font_size=12, italic=True, color=LIGHT_GRAY)
            y += Inches(1.32)

    # Fallback callout
    add_rect(s, Inches(0.3), Inches(6.2), Inches(12.7), Inches(0.95),
             RGBColor(0x1E, 0x14, 0x08))
    add_text(s,
             "Fallback demo: API key removed → system auto-detects, logs WARNING, returns weighted-scorer results labeled [AI unavailable]",
             Inches(0.5), Inches(6.3), Inches(12.3), Inches(0.75),
             font_size=15, color=GOLD, align=PP_ALIGN.CENTER)


def slide_testing(prs):
    """Slide 6 — Testing & Reliability"""
    s = blank_slide(prs)
    bg(s)
    accent_bar(s, color=GOLD)

    add_text(s, "Testing & Reliability — 4 Mechanisms",
             Inches(0.5), Inches(0.2), Inches(12), Inches(0.7),
             font_size=34, bold=True, color=WHITE)

    # Test results block
    add_rect(s, Inches(0.3), Inches(1.0), Inches(5.8), Inches(5.6),
             RGBColor(0x0A, 0x0A, 0x18))

    test_lines = [
        ("$ pytest tests/ -v", 13, ACCENT2, False),
        ("", 8, WHITE, False),
        ("PASSED  test_parse_nl_validates_required_keys",    13, GREEN, False),
        ("PASSED  test_ai_rerank_only_returns_catalog_songs",13, GREEN, False),
        ("PASSED  test_ai_rerank_confidence_scores_returned",13, GREEN, False),
        ("PASSED  test_ai_fallback_on_api_error",            13, GREEN, False),
        ("PASSED  test_ai_recommender_consistency",          13, GREEN, False),
        ("PASSED  test_recommend_returns_songs_sorted",      13, GREEN, False),
        ("PASSED  test_explain_recommendation_non_empty",    13, GREEN, False),
        ("", 8, WHITE, False),
        ("7 passed in 0.28s",                               16, GOLD,  True),
        ("All API calls mocked — no key required",           13, LIGHT_GRAY, False),
    ]

    txb = s.shapes.add_textbox(Inches(0.45), Inches(1.1), Inches(5.5), Inches(5.4))
    tf = txb.text_frame
    tf.word_wrap = False
    first = True
    for txt, sz, col, bold in test_lines:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(2)
        run = p.add_run()
        run.text = txt
        run.font.size = Pt(sz)
        run.font.color.rgb = col
        run.font.bold = bold

    # Right side: 4 mechanisms
    mechanisms = [
        (GREEN,  "1  Automated Tests",
                 "7 pytest tests — guardrail, fallback, consistency,\nconfidence scores. Zero API calls needed."),
        (GOLD,   "2  Confidence Scoring",
                 "Claude self-reports 0.0–1.0 per song.\nAvg confidence: ~0.88 pop/lofi · ~0.65 niche genres."),
        (ACCENT, "3  Logging & Error Handling",
                 "Every Claude call → logs/YYYY-MM-DD.log\nWARNING on hallucination, fallback, omission."),
        (ACCENT2,"4  Human Evaluation",
                 "5 profiles tested manually. 4/5 top results\nmatched expectation. 1 edge-case bias documented."),
    ]

    y = Inches(1.05)
    for col, heading, body in mechanisms:
        add_rect(s, Inches(6.4), y, Inches(0.1), Inches(1.1), col)
        add_rect(s, Inches(6.5), y, Inches(6.5), Inches(1.1), RGBColor(0x16, 0x16, 0x28))
        add_text(s, heading,
                 Inches(6.7), y + Inches(0.08), Inches(6.1), Inches(0.4),
                 font_size=17, bold=True, color=col)
        add_text(s, body,
                 Inches(6.7), y + Inches(0.48), Inches(6.1), Inches(0.55),
                 font_size=13, color=LIGHT_GRAY)
        y += Inches(1.22)

    # Log sample
    add_rect(s, Inches(6.4), y + Inches(0.1), Inches(6.5), Inches(0.85),
             RGBColor(0x0A, 0x0A, 0x18))
    add_text(s,
             "WARNING  vibefinder.ai  AI returned unknown title 'HALLUCINATED TRACK' — skipping\n"
             "WARNING  vibefinder.ai  AI pipeline failed (APIError: rate limit) — falling back",
             Inches(6.5), y + Inches(0.15), Inches(6.3), Inches(0.75),
             font_size=11, color=GOLD)


def slide_learned(prs):
    """Slide 7 — What I Learned"""
    s = blank_slide(prs)
    bg(s)
    accent_bar(s, color=ACCENT)

    add_text(s, "What I Learned",
             Inches(0.5), Inches(0.2), Inches(12), Inches(0.7),
             font_size=34, bold=True, color=WHITE)

    learnings = [
        (ACCENT,   "The AI layer is not the intelligence — the system design is.",
                   "The weighted scorer already produced plausible results. Claude's value was making it "
                   "accessible (natural language in) and more nuanced (musical reasoning at the ranking step). "
                   "The features and weights you choose matter more than which model you call."),
        (ACCENT2,  "Testing LLM systems means testing the plumbing, not the model.",
                   "You mock the AI and verify that your code handles valid output, malformed output, "
                   "and failure correctly. The model's quality is measured through sample interactions "
                   "and confidence scores — not assertions."),
        (GOLD,     "\"Right by the formula\" ≠ \"right for the human.\"",
                   "The Conflicted Raver edge case (EDM genre + melancholy mood) returned Ultraviolet Drop "
                   "because genre + energy outweighed mood. Correct by the weights I wrote — but the human "
                   "wanted something moody. That gap is where AI alignment lives, even in a music app."),
        (GREEN,    "Fallback paths are not optional — they are the product.",
                   "A system that goes down when the API is unavailable is not a system, it is a demo. "
                   "Building the weighted scorer as a standalone fallback meant the user always gets "
                   "recommendations regardless of Claude's availability."),
    ]

    y = Inches(1.05)
    for col, heading, body in learnings:
        add_rect(s, Inches(0.3), y, Inches(0.12), Inches(1.45), col)
        add_rect(s, Inches(0.45), y, Inches(12.5), Inches(1.45), RGBColor(0x14, 0x14, 0x26))
        add_text(s, heading,
                 Inches(0.65), y + Inches(0.1), Inches(12.1), Inches(0.45),
                 font_size=16, bold=True, color=col)
        add_text(s, body,
                 Inches(0.65), y + Inches(0.55), Inches(12.1), Inches(0.8),
                 font_size=14, color=LIGHT_GRAY)
        y += Inches(1.56)


def slide_reflection(prs):
    """Slide 8 — Reflection & GitHub"""
    s = blank_slide(prs)
    bg(s)
    # Side accent panel
    add_rect(s, SLIDE_W - Inches(3.8), 0, Inches(3.8), SLIDE_H,
             RGBColor(0x10, 0x08, 0x28))
    accent_bar(s, y=0, color=ACCENT2)
    accent_bar(s, y=SLIDE_H - Inches(0.06), color=ACCENT)

    add_text(s, "Reflection",
             Inches(0.5), Inches(0.25), Inches(9.0), Inches(0.7),
             font_size=38, bold=True, color=WHITE)

    add_text(s, "What this project says about me as an AI engineer",
             Inches(0.5), Inches(0.95), Inches(9.0), Inches(0.5),
             font_size=20, italic=True, color=ACCENT2)

    reflection = (
        "VibeFinder AI shows that I think about AI as a system design problem, not just an API call. "
        "I extended a working baseline rather than replacing it — the weighted scorer stayed intact "
        "because it was correct and transparent, and Claude was layered on top to add natural language "
        "access and musical nuance. I built for failure from the start: the fallback path, the output "
        "guardrail, the structured JSON schema, and the logging were not afterthoughts — they were the "
        "design. The confidence scoring shows I think about measuring reliability, not just assuming it. "
        "And the 7-test suite, run without any API key, shows I build things that other people can verify.\n\n"
        "The \"Conflicted Raver\" edge case was the most honest moment of the project: the system did "
        "exactly what I told it to do, and the result felt wrong. That gap — between a formula that is "
        "correct and an output that is right for the human — is the problem I want to keep working on."
    )

    add_rect(s, Inches(0.3), Inches(1.55), Inches(9.1), Inches(4.3),
             RGBColor(0x12, 0x12, 0x24))
    txb = s.shapes.add_textbox(Inches(0.5), Inches(1.65), Inches(8.7), Inches(4.1))
    tf = txb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = reflection
    run.font.size = Pt(15.5)
    run.font.color.rgb = LIGHT_GRAY

    # Right panel — GitHub + tech
    add_text(s, "GitHub", SLIDE_W - Inches(3.5), Inches(0.6), Inches(3.2), Inches(0.5),
             font_size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    add_rect(s, SLIDE_W - Inches(3.6), Inches(1.1), Inches(3.35), Inches(0.04), ACCENT)

    add_text(s, "github.com/guka199/\napplied-music-\nrecommendations",
             SLIDE_W - Inches(3.5), Inches(1.25), Inches(3.2), Inches(1.1),
             font_size=14, color=ACCENT2, align=PP_ALIGN.CENTER)

    add_text(s, "Tech Stack", SLIDE_W - Inches(3.5), Inches(2.55), Inches(3.2), Inches(0.4),
             font_size=17, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    add_rect(s, SLIDE_W - Inches(3.6), Inches(2.95), Inches(3.35), Inches(0.04), ACCENT2)

    stack = [
        "Python 3.14",
        "Anthropic SDK",
        "claude-opus-4-7",
        "json_schema output_config",
        "cache_control: ephemeral",
        "pytest + unittest.mock",
        "python-dotenv",
        "python-pptx",
    ]
    txb2 = s.shapes.add_textbox(SLIDE_W - Inches(3.5), Inches(3.1), Inches(3.2), Inches(3.8))
    tf2 = txb2.text_frame
    tf2.word_wrap = True
    first = True
    for item in stack:
        if first:
            p2 = tf2.paragraphs[0]
            first = False
        else:
            p2 = tf2.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        p2.space_before = Pt(5)
        run2 = p2.add_run()
        run2.text = item
        run2.font.size = Pt(14)
        run2.font.color.rgb = LIGHT_GRAY

    # Bottom bar
    add_rect(s, 0, SLIDE_H - Inches(0.55), SLIDE_W, Inches(0.49),
             RGBColor(0x16, 0x08, 0x30))
    add_text(s, "Guram Janashia  ·  AI 110  ·  April 2026  ·  github.com/guka199/applied-music-recommendations",
             Inches(0.3), SLIDE_H - Inches(0.52), SLIDE_W - Inches(0.6), Inches(0.44),
             font_size=13, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)


# ── Main ─────────────────────────────────────────────────────────────────────

def main():
    prs = new_prs()

    slide_title(prs)
    slide_origin(prs)
    slide_what_i_built(prs)
    slide_architecture(prs)
    slide_demo(prs)
    slide_testing(prs)
    slide_learned(prs)
    slide_reflection(prs)

    out = "VibeFinder_AI_Presentation.pptx"
    prs.save(out)
    print(f"Saved: {out}  ({prs.slides.__len__()} slides)")


if __name__ == "__main__":
    main()
